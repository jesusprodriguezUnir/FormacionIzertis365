"""
Management command para extraer slides de archivos PowerPoint reales con imágenes.
"""
import os
import re
import shutil
import subprocess
import tempfile
from django.core.management.base import BaseCommand
from django.core.files.base import ContentFile
from pptx import Presentation
from core.models import Resource, Slide


class Command(BaseCommand):
    help = 'Extrae slides de archivos PowerPoint con imágenes y texto'

    def _find_soffice(self):
        soffice = shutil.which('soffice') or shutil.which('soffice.exe')
        if soffice:
            return soffice

        possible_paths = [
            r"C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            r"C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None

    def _extract_slide_images(self, pptx_path, output_dir):
        soffice_path = self._find_soffice()
        if not soffice_path:
            return None, 'LibreOffice (soffice) no está instalado o no está en PATH.'

        # 1) Convertir a PDF con LibreOffice
        cmd = [
            soffice_path,
            '--headless',
            '--convert-to',
            'pdf',
            '--outdir',
            output_dir,
            pptx_path,
        ]

        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            error_message = (result.stderr or result.stdout or '').strip()
            return None, f'Error al convertir PPTX a PDF: {error_message}'

        # Buscar el PDF generado
        pdf_files = [
            os.path.join(output_dir, f)
            for f in os.listdir(output_dir)
            if f.lower().endswith('.pdf')
        ]
        if not pdf_files:
            return None, 'No se generó PDF durante la conversión.'

        pdf_path = pdf_files[0]

        # 2) Convertir PDF a PNG usando PyMuPDF (fitz)
        try:
            import fitz  # PyMuPDF
        except Exception:
            return None, 'PyMuPDF (fitz) no está instalado. Instálalo con: pip install pymupdf'

        image_files = []
        try:
            doc = fitz.open(pdf_path)
            for page_index in range(len(doc)):
                page = doc.load_page(page_index)
                pix = page.get_pixmap(dpi=200)
                output_path = os.path.join(output_dir, f"slide_{page_index + 1}.png")
                pix.save(output_path)
                image_files.append(output_path)
        finally:
            doc.close()

        return image_files, None

    def add_arguments(self, parser):
        parser.add_argument(
            '--pptx-path',
            type=str,
            required=True,
            help='Ruta al archivo PowerPoint'
        )
        parser.add_argument(
            '--resource-slug',
            type=str,
            required=True,
            help='Slug del recurso al que asociar los slides'
        )

    def handle(self, *args, **options):
        pptx_path = options['pptx_path']
        resource_slug = options['resource_slug']
        
        # Verificar que existe el archivo
        if not os.path.exists(pptx_path):
            self.stdout.write(self.style.ERROR(f'✗ No se encontró el archivo: {pptx_path}'))
            return
        
        # Buscar el recurso
        try:
            resource = Resource.objects.get(slug=resource_slug)
        except Resource.DoesNotExist:
            self.stdout.write(self.style.ERROR(f'✗ No se encontró el recurso con slug: {resource_slug}'))
            return
        
        self.stdout.write(self.style.SUCCESS(f'✓ Recurso encontrado: {resource.title}'))
        self.stdout.write(self.style.WARNING(f'Procesando: {pptx_path}'))
        
        # Eliminar slides existentes
        deleted_count = Slide.objects.filter(resource=resource).count()
        if deleted_count > 0:
            Slide.objects.filter(resource=resource).delete()
            self.stdout.write(self.style.WARNING(f'✓ Eliminados {deleted_count} slides existentes'))
        
        # Cargar la presentación
        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            self.stdout.write(self.style.ERROR(f'✗ Error al abrir PowerPoint: {e}'))
            return

        # Convertir slides a imágenes (si LibreOffice está disponible)
        image_files = []
        image_error = None
        with tempfile.TemporaryDirectory() as temp_dir:
            image_files, image_error = self._extract_slide_images(pptx_path, temp_dir)
            if image_error:
                self.stdout.write(self.style.WARNING(f'⚠ {image_error}'))
                self.stdout.write(self.style.WARNING('Se continuará con contenido de texto.'))
                image_files = []
            else:
                self.stdout.write(self.style.SUCCESS(f'✓ Imágenes generadas: {len(image_files)}'))

            created_count = 0

            # Procesar cada slide
            for idx, slide in enumerate(prs.slides, start=1):
                # Extraer texto
                text_content = []
                title_text = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if not title_text and len(text) < 200:  # Primer texto corto como título
                            title_text = text
                        text_content.append(text)
                
                # Combinar todo el texto
                full_text = "\n\n".join(text_content)
                
                # Si no hay título, usar primeras palabras
                if not title_text and full_text:
                    title_text = full_text.split('\n')[0][:200]
                
                if not title_text:
                    title_text = f"Slide {idx}"
                
                # Crear el slide
                slide_obj = Slide.objects.create(
                    resource=resource,
                    slide_number=idx,
                    title=title_text,
                    content=full_text,
                    order=idx
                )

                # Asignar imagen si existe
                if image_files and idx <= len(image_files):
                    image_path = image_files[idx - 1]
                    if os.path.exists(image_path):
                        with open(image_path, 'rb') as img_file:
                            filename = f"{resource.slug}-slide-{idx}.png"
                            slide_obj.image.save(filename, ContentFile(img_file.read()), save=True)
                
                created_count += 1
                self.stdout.write(f'  ✓ Slide {idx}: {title_text[:50]}...')
            self.stdout.write(self.style.SUCCESS(f'\n✓ {created_count} slides creados correctamente'))
